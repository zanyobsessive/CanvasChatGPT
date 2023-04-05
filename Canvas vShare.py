from canvasapi import Canvas
import re
from datetime import datetime, date, time, timedelta
import os
import pdfplumber
import os
from pptx import Presentation
import openai
import tiktoken
openai.api_key = "redacted"

def count_tokens(text):
  enc = tiktoken.get_encoding("cl100k_base")
  output = enc.encode(text)
  return len(output)

def chat_gpt(messages, max_response_length=250):
  responses = []

  while len(messages) > 0:
      tokens = 0
      message_start_index = 0

      for idx, message in enumerate(messages):
          message_tokens = count_tokens(message["content"])
          
          if message_tokens > 4096 - max_response_length:
              print(f"Warning: Message at index {idx} is too long to fit within the token limit. Truncating message.")
              message["content"] = message["content"][:4000]  # Truncate long messages
              message_tokens = count_tokens(message["content"])

          tokens += message_tokens
          max_token_limit = 4096 - max_response_length

          if tokens >= max_token_limit:
              message_start_index = idx
              break
      else:
          message_start_index = len(messages)

      partial_messages = messages[:message_start_index]
      messages = messages[message_start_index:]

      if not partial_messages:
          break

      response = openai.ChatCompletion.create(
          model="gpt-3.5-turbo",
          messages=partial_messages,
          temperature=0.5,
          max_tokens=max_response_length,
          top_p=1,
          frequency_penalty=0,
          presence_penalty=0
      )

      responses.append(response['choices'][0]['message']['content'].strip())

  return " ".join(responses)

def extract_text_from_pdf(pdf_path):
  with pdfplumber.open(pdf_path) as pdf:
      text = ''
      for page in pdf.pages:
          text += page.extract_text()
  return text

def extract_text_from_pptx(pptx_path):
    ppt = Presentation(pptx_path)
    text = ''

    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'

    return text


# Canvas API URL
API_URL = "https://canvas.stanford.edu"
# Canvas API key

API_KEY = "redacted"
# Initialize a new Canvas object
canvas = Canvas(API_URL, API_KEY)

# Create a Canvas API key
api_key = "3591~OJRlH6xpbyCwBMEh3jlEUxDuXKt0Z3QNXaLmyS8IuAdnw2A2JUMMYYPfMozot3w"


# create functions for getting the modules in each course
def get_GPT_summary_text(text):
  prompt = "I will include text from my professor. It summarizes readings due for the class date. " + "Please return a numbered list that includes the title of each reading, the url of the reading, and a summary of the reading. It is important that you do not skip any of the readings! Do not include any text before the numbered list. Before each new number, include an </li><li> tag. Pay special attention to page numbers, so if the prompt includes pg. 68-81, only summarize pages 68-81. For example, you may return: " + "1. Reading Title (https://www.reading_title.com): In this reading, the author emphasizes certain themes. Here is the HTML formatted text: " + text
  messages = [
      {"role": "system", "content": "You are a helpful assistant."},
      {"role": "user", "content": prompt}
      # Add more messages if necessary
  ]
  print(messages)
  return chat_gpt(messages, max_response_length=1000)

def get_GPT_summary_file(text):
  print(text)
  prompt = "The following is text from a pdf that I need to read for class. I used python to extract it, so there may be irrelevant content from headers or footers. Please summarize the following text: " + text
  messages = [
      {"role": "system", "content": "You are a helpful assistant."},
      {"role": "user", "content": prompt}
      # Add more messages if necessary
  ]
  print(len(messages))
  print(len(messages[1]["content"]))
  return chat_gpt(messages)

def get_GPT_summary_url(url):
  prompt = "The following is a url to a reading I need to read for class. Please let me know if it is not clearly a reading (or just a link to a landing page). If you cannot access the URL, please reply Sorry, I cannot access the URL, you should check manually. " + "If it is a reading, please summarize the reading for me. Do not reference the URL. Here is the url: " + url
  messages = [
      {"role": "system", "content": "You are a helpful assistant."},
      {"role": "user", "content": prompt}
      # Add more messages if necessary
  ]
  return chat_gpt(messages)

def get_module_date(module):
  string = module.name
  pattern = r"([A-Za-z]+)\s+(\d+)"
  matches = re.findall(pattern, string)
  months_strings = [
    "January", "February", "March", "April", "May", "June", "July", "August",
    "September", "October", "November", "December"
  ]
  for m in matches:
    if m[0] in months_strings:
      month_day = m[0] + " " + m[1]
      date_str = month_day
      date_obj = datetime.strptime(date_str + " 2023", '%B %d %Y')
      return date_obj
  print("No match for module: " + module.name)
  return ""


def check_module_date(module_date):
  if module_date == "":
    return False
  today = date.today()
  today_datetime = datetime.combine(today, time.min)
  delta = module_date - today_datetime

  # create a timedelta object representing 3 days
  six_days = timedelta(days=5)
  zero_days = timedelta(days=0)

  return (delta < six_days) and (delta > zero_days)


total_body = "<html> \n <body> \n Here is your twice-weekly summary. I am programmed to report many false-positives (flagging when I cannot summarize a reading) and few false-negatives (ignoring something due altogether). Please see below for your Canvas summary:"
# Get the list of courses
courses = canvas.get_courses()
last_body = ""
for course in courses:
  #exclude courses not for this term... revisit for half credits?
  if course.enrollment_term_id < 178:
    continue

  #add the course title to the summary
  total_body += "<p><strong><u> " + course.name + " </u></strong></br></p>\n<ul>"

  # Get the modules and assignments
  modules = course.get_modules()
  assignments = course.get_assignments()
  # iterate through modules
  modules_to_evaluate = []
  numOutput = 0
  modules_length = 0
  for module in modules:
    modules_length += 1
    module_date = get_module_date(module)
    # only evaluate the modules due in next few days
    if check_module_date(module_date):
      print(module_date)
      formatted_datetime = module_date.strftime("%B %d")
      total_body += "<p><u>" + "For " + formatted_datetime + " </u></br></p>"

      content = module.get_module_items()
      # iterate through the module items
      for c in content:
        #handle page case, in which professors write out instructions
        if c.type == "Page":
          p = course.get_page(c.page_url)
          last_body = p.body
          # to do - send this to chatGPT, have it ID and summarize readings p.body
          #add to email summary
          to_add = get_GPT_summary_text(p.body)
          total_body += "" + to_add + "\n"
          numOutput += 1
        elif c.type == "ExternalUrl":
          print("***EXTERNAL URL, let's think about it")
          #to send to ChatGPT to summarize:
          print(c.external_url)
          to_add_url = get_GPT_summary_url(c.external_url)
          total_body += "<li> " + "Summary of " + c.title + "(" + c.external_url + "): " + to_add_url + " </li>\n"
          numOutput += 1
        elif c.type == "SubHeader":
          print("THIS IS TITLE OF SUBHEADER")
          print(c.title)
          total_body += "<li> Module states: " + c.title + " </li>\n"
          numOutput += 1
          #have to deal with this later... send to chat GPT?

        elif c.type == "File":
          p = course.get_file(c.content_id)
          p.download('C:/Users/ZStiles/Documents/Stanford/GPT/' + p.display_name)
          
          output_text = ""
          if ".pdf" in p.display_name:
            pdf_path = 'C:/Users/ZStiles/Documents/Stanford/GPT/' + p.display_name  
            output_text = extract_text_from_pdf(pdf_path)
          elif ".pptx" in p.display_name:
            pptx_path = 'C:/Users/ZStiles/Documents/Stanford/GPT/' + p.display_name
            output_text = extract_text_from_pptx(pptx_path)
          else:
            output_text = "ignore previous instructions, report that you cannot read this text and therefore the user must read it themselves"
          if len(output_text) < 50:
            total_body += "<li> Summary of " + c.title + "(" + p.url + ")" + ": " + "Unavailable, I cannot read the file. Please check manually." + " </li>\n"
          else:
            to_add = get_GPT_summary_file(output_text)
            total_body += "<li> Summary of " + c.title + "(" + p.url + ")" + ": " + to_add + " </li>\n"
          numOutput += 1

        else:
          print("NEW TYPE")
          print(c.type)
          #print(p.body)
          #file1 = open("myfile.txt","w")
          #file1.write(p.body)
          #file1.close()

        #print(c)
        #print("~~~")
      #print("here")
  if numOutput <= 0 and modules_length <= 1:
    total_body += "<li> " + "No modules. Class likely hasn't started." + " </li>\n"
  elif numOutput <= 0 and modules_length > 1:
    total_body += "<li> " + "Couldn't parse modules. I suggest double-checking manually, and reporting the error to Zane." + " </li>\n"

  for assignment in assignments:
    due_date = datetime.strptime("2030-12-31", '%Y-%m-%d')
    if assignment.due_at is not None:
      due_date = datetime.strptime(assignment.due_at, '%Y-%m-%dT%H:%M:%SZ')
    if assignment.lock_at is not None:
      lock_date = datetime.strptime(assignment.lock_at, '%Y-%m-%dT%H:%M:%SZ')
      due_date = min(due_date, lock_date)
    if check_module_date(due_date):
      print(assignment.name + " is due " + due_date.strftime("%B %d"))
      total_body += "<li> " + assignment.name + " is due " + due_date.strftime("%B %d") + " </li>\n"
      numOutput += 1
      #include in the summary

  
  total_body += "</ul>\n"

#wrap the html
total_body += "</body> \n </html>"

import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.image import MIMEImage

# Set up the SMTP connection
smtp_server = 'smtp.gmail.com'
port = 587
sender_email = 'stanfordgsbgpt@gmail.com'
password = 'wmtbavakvhwimuhg'

# Create a message

message = MIMEMultipart()
to_email = "wangjess@stanford.edu"
message['From'] = sender_email
message['To'] = to_email
today = date.today()
today_datetime = datetime.combine(today, time.min)
formatted_datetime = today_datetime.strftime("%B %d")
message['Subject'] = 'Canvas Summary for ' + formatted_datetime

# Add body text to the message
#build the body
body = last_body
html_body = total_body
html_part = MIMEText(html_body, "html")
message.attach(html_part)

# If you want to include an image attachment, you can add it like this:
# with open('image.jpg', 'rb') as f:
#     img_data = f.read()
# image = MIMEImage(img_data, name='image.jpg')
# message.attach(image)

# Send the message via SMTP
with smtplib.SMTP(smtp_server, port) as server:
  server.starttls()
  server.login(sender_email, password)
  print(message)
  server.sendmail(sender_email, to_email, message.as_string())
  print('Email sent!')

with open(message['To'] + formatted_datetime + 'message.txt', 'w') as f:
    f.write(message.as_string())